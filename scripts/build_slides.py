"""
build_slides.py — Vaartha Presentation Builder
Generates a 15-slide .pptx with dark theme, embedded charts, and talking points.
Upload to Google Drive → it auto-converts to Google Slides.

Usage:
    python scripts/build_slides.py
Output:
    outputs/Vaartha_Presentation.pptx
"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "src"))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util

from config import CHARTS, OUTPUTS

# ── Palette ────────────────────────────────────────────────────────────────
BG          = RGBColor(0x0A, 0x0A, 0x0F)
SURFACE     = RGBColor(0x1E, 0x1E, 0x2E)
TEXT_MAIN   = RGBColor(0xE2, 0xE8, 0xF0)
TEXT_DIM    = RGBColor(0x94, 0xA3, 0xB8)
ST1_COLOR   = RGBColor(0x00, 0xFF, 0xB2)
ST2_COLOR   = RGBColor(0xFF, 0x6B, 0x35)
ST3_COLOR   = RGBColor(0xA7, 0x8B, 0xFA)
ST4_COLOR   = RGBColor(0xFA, 0xCC, 0x15)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_PINK = RGBColor(0xF4, 0x72, 0xB6)

# Slide dimensions: widescreen 16:9
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    """Create a blank widescreen presentation."""
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    """Add a completely blank slide."""
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def fill_bg(slide, color=BG):
    """Fill slide background with solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=24, bold=False, color=TEXT_MAIN,
             align=PP_ALIGN.LEFT, wrap=True):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Inter"
    return txBox


def add_bullet_box(slide, bullets, left, top, width, height,
                   font_size=16, color=TEXT_DIM, marker_color=None):
    """Add a text box with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"  {bullet}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Inter"
    return txBox


def add_image(slide, img_path, left, top, width, height):
    """Add an image if it exists and is a supported format; otherwise placeholder."""
    path = Path(img_path)
    supported = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff'}
    if path.exists() and path.suffix.lower() in supported:
        slide.shapes.add_picture(str(path), left, top, width, height)
    else:
        add_rect(slide, left, top, width, height, SURFACE)
        note = f"[Interactive chart: {path.name}]\nOpen in browser for full interactivity."
        add_text(slide, note,
                 left + Inches(0.2), top + height / 2 - Inches(0.4),
                 width - Inches(0.4), Inches(0.8),
                 font_size=14, color=TEXT_DIM, align=PP_ALIGN.CENTER)


def pill_label(slide, text, left, top, color):
    """Add a small colored pill label (e.g. 'ST1')."""
    w, h = Inches(1.1), Inches(0.35)
    rect = add_rect(slide, left, top, w, h, color)
    add_text(slide, text,
             left, top + Pt(2), w, h,
             font_size=13, bold=True,
             color=RGBColor(0x0A, 0x0A, 0x0F),
             align=PP_ALIGN.CENTER)


def speaker_tag(slide, name):
    """Small speaker attribution in bottom-right."""
    add_text(slide, f"▶  {name}",
             Inches(11.0), Inches(7.1),
             Inches(2.2), Inches(0.35),
             font_size=11, color=TEXT_DIM,
             align=PP_ALIGN.RIGHT)


# ── SECTION HEADER helper ──────────────────────────────────────────────────
def section_bar(slide, label, accent):
    """Thin accent bar at top of slide with section label."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.07), accent)
    add_text(slide, label,
             Inches(0.4), Inches(0.08),
             Inches(4), Inches(0.35),
             font_size=12, bold=True, color=accent)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════

def slide_01_hook(prs):
    """Slide 1 — Hook"""
    sl = blank_slide(prs)
    fill_bg(sl)

    # Big background text (decorative)
    add_text(sl, "?", Inches(8.5), Inches(0.5), Inches(4), Inches(6.5),
             font_size=300, bold=True,
             color=RGBColor(0x1E, 0x1E, 0x2E), align=PP_ALIGN.CENTER)

    add_text(sl, "What happens when a country\nweaponizes a mineral?",
             Inches(0.5), Inches(1.2), Inches(7.5), Inches(2.5),
             font_size=40, bold=True, color=WHITE)

    add_text(sl,
             "In 2023, China banned gallium & germanium exports overnight.\n"
             "NVIDIA stock dropped. Intel announced a $20B fab. Congress passed the CHIPS Act.\n\n"
             "Is this pattern measurable? We built a 13-year data pipeline to find out.",
             Inches(0.5), Inches(3.9), Inches(7.5), Inches(2.2),
             font_size=17, color=TEXT_DIM)

    # Stat callout
    add_rect(sl, Inches(0.5), Inches(6.3), Inches(3.0), Inches(0.75), SURFACE)
    add_text(sl, "China controls 77% of global gallium",
             Inches(0.6), Inches(6.35), Inches(2.9), Inches(0.65),
             font_size=14, bold=True, color=ST1_COLOR)

    speaker_tag(sl, "Tarun")
    return sl


def slide_02_thesis(prs):
    """Slide 2 — Thesis & Causal Chain"""
    sl = blank_slide(prs)
    fill_bg(sl)

    add_text(sl, "The Causal Chain We're Testing",
             Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
             font_size=32, bold=True, color=WHITE)

    add_text(sl, "Thesis: Geopolitical fragmentation distorts critical mineral supply chains, "
             "forcing corporations to aggressively restructure CapEx and R&D to survive.",
             Inches(0.5), Inches(1.1), Inches(12), Inches(0.7),
             font_size=16, color=TEXT_DIM)

    # Chain boxes
    boxes = [
        ("ST2", "Geopolitical\nShock", ST2_COLOR, Inches(0.4)),
        ("ST1", "Supply Chain\nDistortion", ST1_COLOR, Inches(3.3)),
        ("ST3", "Energy Policy\nAmplifies", ST3_COLOR, Inches(6.2)),
        ("ST4", "Corporate\nCapEx Response", ST4_COLOR, Inches(9.1)),
    ]
    for tag, label, color, left in boxes:
        add_rect(sl, left, Inches(2.2), Inches(2.7), Inches(2.2), SURFACE)
        add_rect(sl, left, Inches(2.2), Inches(2.7), Inches(0.08), color)
        add_text(sl, tag, left + Inches(0.1), Inches(2.3),
                 Inches(1), Inches(0.4),
                 font_size=13, bold=True, color=color)
        add_text(sl, label, left + Inches(0.1), Inches(2.75),
                 Inches(2.5), Inches(1.5),
                 font_size=20, bold=True, color=WHITE)

    # Arrows between boxes
    for ax in [Inches(3.1), Inches(6.0), Inches(8.9)]:
        add_text(sl, "→", ax, Inches(2.9), Inches(0.2), Inches(0.5),
                 font_size=28, bold=True, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    add_text(sl, "9 datasets  ·  121 countries  ·  10 companies  ·  13 years (2010–2022)",
             Inches(0.5), Inches(6.6), Inches(12), Inches(0.5),
             font_size=15, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    speaker_tag(sl, "Tarun")
    return sl


def slide_03_st1_hhi(prs):
    """Slide 3 — ST1: HHI Over Time"""
    sl = blank_slide(prs)
    fill_bg(sl)
    section_bar(sl, "ST1 — Critical Minerals", ST1_COLOR)

    add_text(sl, "A Handful of Countries Control Everything",
             Inches(0.4), Inches(0.5), Inches(8), Inches(0.7),
             font_size=28, bold=True, color=WHITE)

    add_image(sl, CHARTS / "st1_hhi_over_time.png",
              Inches(0.4), Inches(1.3), Inches(8.0), Inches(5.5))

    # Stat callouts on right
    stats = [
        ("Gallium", "HHI 0.77", ST1_COLOR),
        ("Germanium", "HHI 0.76", ST1_COLOR),
        ("Rare Earths", "HHI 0.71", ST1_COLOR),
    ]
    for i, (label, val, col) in enumerate(stats):
        top = Inches(1.5 + i * 1.5)
        add_rect(sl, Inches(8.8), top, Inches(4.0), Inches(1.1), SURFACE)
        add_text(sl, val, Inches(8.9), top + Inches(0.05),
                 Inches(3.8), Inches(0.55),
                 font_size=30, bold=True, color=col)
        add_text(sl, label, Inches(8.9), top + Inches(0.6),
                 Inches(3.8), Inches(0.4),
                 font_size=15, color=TEXT_DIM)

    add_text(sl, "All three dominated by China. Structural, not cyclical.",
             Inches(8.8), Inches(6.2), Inches(4.3), Inches(0.8),
             font_size=14, color=ACCENT_PINK)

    speaker_tag(sl, "Tarun")
    return sl


def slide_04_st1_sankey(prs):
    """Slide 4 — ST1: Trade Flows"""
    sl = blank_slide(prs)
    fill_bg(sl)
    section_bar(sl, "ST1 — Critical Minerals", ST1_COLOR)

    add_text(sl, "One Disruption = Multiple Simultaneous Shocks",
             Inches(0.4), Inches(0.5), Inches(12), Inches(0.7),
             font_size=28, bold=True, color=WHITE)

    # Placeholder for Sankey (HTML — show screenshot or describe)
    add_image(sl, CHARTS / "st1_sankey_flows.html",
              Inches(0.4), Inches(1.3), Inches(8.5), Inches(5.3))

    add_text(sl, "The compounded\nexposure problem:",
             Inches(9.2), Inches(1.5), Inches(3.8), Inches(0.9),
             font_size=17, bold=True, color=WHITE)

    bullets = [
        "China appears in gallium, germanium, rare earths & graphite flows simultaneously",
        "One diplomatic incident = 4 correlated supply shocks",
        "No diversification alternative exists at scale today",
    ]
    add_bullet_box(sl, bullets,
                   Inches(9.2), Inches(2.6), Inches(3.8), Inches(3.5),
                   font_size=15, color=TEXT_DIM)

    speaker_tag(sl, "Tarun")
    return sl


def slide_05_st1_restrictions(prs):
    """Slide 5 — ST1: OECD Restrictions"""
    sl = blank_slide(prs)
    fill_bg(sl)
    section_bar(sl, "ST1 — Critical Minerals", ST1_COLOR)

    add_text(sl, "Governments Are Weaponizing Export Restrictions",
             Inches(0.4), Inches(0.5), Inches(12), Inches(0.7),
             font_size=28, bold=True, color=WHITE)

    add_image(sl, CHARTS / "st1_oecd_restrictions.png",
              Inches(0.4), Inches(1.3), Inches(8.0), Inches(5.5))

    add_text(sl, "Rising every year\nsince 2018",
             Inches(8.8), Inches(1.5), Inches(4.0), Inches(0.9),
             font_size=22, bold=True, color=ST2_COLOR)

    bullets = [
        "2022 spike aligns exactly with Ukraine invasion",
        "Mineral access is now a geopolitical instrument",
        "This links ST1 directly to ST2 — which Chaitanya will cover",
    ]
    add_bullet_box(sl, bullets,
                   Inches(8.8), Inches(2.6), Inches(4.0), Inches(3.0),
                   font_size=15, color=TEXT_DIM)

    speaker_tag(sl, "Tarun → handoff to Chaitanya")
    return sl


def slide_06_st2_gpr(prs):
    """Slide 6 — ST2: GPR Timeseries"""
    sl = blank_slide(prs)
    fill_bg(sl)
    section_bar(sl, "ST2 — Geopolitical Risk", ST2_COLOR)

    add_text(sl, "Geopolitical Risk Is Measurable — and It Transmits",
             Inches(0.4), Inches(0.5), Inches(12), Inches(0.7),
             font_size=28, bold=True, color=WHITE)

    add_image(sl, CHARTS / "st2_gpr_timeseries.png",
              Inches(0.4), Inches(1.3), Inches(8.0), Inches(5.5))

    bullets = [
        "GPR built from newspaper article counts — 2,000+ sources",
        "GPRT = threats (forward-looking)",
        "GPRA = acts (realized events) — produces faster price responses",
        "Every major event is visible: 2014, 2018, 2020, 2022",
    ]
    add_bullet_box(sl, bullets,
                   Inches(8.8), Inches(1.5), Inches(4.0), Inches(4.5),
                   font_size=15, color=TEXT_DIM)

    speaker_tag(sl, "Chaitanya")
    return sl


def slide_07_st2_transmission(prs):
    """Slide 7 — ST2: GPR → GSCPI"""
    sl = blank_slide(prs)
    fill_bg(sl)
    section_bar(sl, "ST2 — Geopolitical Risk", ST2_COLOR)

    add_text(sl, "GPR Leads Supply Chain Pressure by 1–2 Months",
             Inches(0.4), Inches(0.5), Inches(12), Inches(0.7),
             font_size=28, bold=True, color=WHITE)

    add_image(sl, CHARTS / "st2_gpr_gscpi_correlation.png",
              Inches(0.4), Inches(1.3), Inches(8.0), Inches(5.5))

    add_rect(sl, Inches(8.8), Inches(1.5), Inches(4.0), Inches(1.0), SURFACE)
    add_text(sl, "Transmission is\ngetting stronger",
             Inches(8.9), Inches(1.55), Inches(3.8), Inches(0.9),
             font_size=20, bold=True, color=ST2_COLOR)

    bullets = [
        "Positive correlation persistent since 2016",
        "Strengthened sharply after 2020",
        "Geopolitical risk now a reliable predictor of supply chain disruption",
        "The relationship is structural, not episodic",
    ]
    add_bullet_box(sl, bullets,
                   Inches(8.8), Inches(2.7), Inches(4.0), Inches(3.5),
                   font_size=15, color=TEXT_DIM)

    speaker_tag(sl, "Chaitanya")
    return sl


def slide_08_st2_ukraine(prs):
    """Slide 8 — ST2: Ukraine Event Study"""
    sl = blank_slide(prs)
    fill_bg(sl)
    section_bar(sl, "ST2 — Geopolitical Risk", ST2_COLOR)

    add_text(sl, "The Ukraine Invasion: A Live Test of the Chain",
             Inches(0.4), Inches(0.5), Inches(12), Inches(0.7),
             font_size=28, bold=True, color=WHITE)

    add_image(sl, CHARTS / "st2_event_study_ukraine.png",
              Inches(0.4), Inches(1.3), Inches(7.8), Inches(5.5))

    add_text(sl, "The sequence:",
             Inches(8.8), Inches(1.5), Inches(4.0), Inches(0.5),
             font_size=17, bold=True, color=WHITE)

    sequence = [
        "① GPR spikes immediately (Feb 24)",
        "② GSCPI follows — 1–2 months later",
        "③ XLE (Energy) +40%",
        "③ SOXX (Semiconductors) −20%",
        "Exactly what the thesis predicts.",
    ]
    add_bullet_box(sl, sequence,
                   Inches(8.8), Inches(2.1), Inches(4.0), Inches(4.0),
                   font_size=15, color=TEXT_DIM)

    speaker_tag(sl, "Chaitanya → handoff to Tarun")
    return sl


def slide_09_st3_mix(prs):
    """Slide 9 — ST3: Energy Mix"""
    sl = blank_slide(prs)
    fill_bg(sl)
    section_bar(sl, "ST3 — Energy Policy", ST3_COLOR)

    add_text(sl, "Energy Policy Divergence Is Amplifying Demand",
             Inches(0.4), Inches(0.5), Inches(12), Inches(0.7),
             font_size=28, bold=True, color=WHITE)

    add_image(sl, CHARTS / "st3_energy_mix_faceted.png",
              Inches(0.4), Inches(1.3), Inches(8.0), Inches(5.5))

    bullets = [
        "Germany & UK: steep renewables acceleration since 2015",
        "India & Brazil: still predominantly fossil-fuel",
        "This creates structurally different mineral demand trajectories",
        "Policy choices in capitals compound the supply crisis in ST1",
    ]
    add_bullet_box(sl, bullets,
                   Inches(8.8), Inches(1.5), Inches(4.0), Inches(4.5),
                   font_size=15, color=TEXT_DIM)

    speaker_tag(sl, "Tarun")
    return sl


def slide_10_st3_demand(prs):
    """Slide 10 — ST3: Mineral Demand Projection"""
    sl = blank_slide(prs)
    fill_bg(sl)
    section_bar(sl, "ST3 — Energy Policy", ST3_COLOR)

    add_text(sl, "Renewables Growth = Exponential Mineral Demand",
             Inches(0.4), Inches(0.5), Inches(12), Inches(0.7),
             font_size=28, bold=True, color=WHITE)

    add_image(sl, CHARTS / "st3_mineral_demand_projection.png",
              Inches(0.4), Inches(1.3), Inches(8.0), Inches(5.5))

    add_rect(sl, Inches(8.8), Inches(1.5), Inches(4.0), Inches(1.0), SURFACE)
    add_text(sl, "Solar grew 30×\nin 12 years",
             Inches(8.9), Inches(1.55), Inches(3.8), Inches(0.9),
             font_size=22, bold=True, color=ST3_COLOR)

    bullets = [
        "IEA 2023 intensity factors: 3t Si/MW solar; 2.5t Cu/MW wind",
        "Implied silicon demand tripled in the last 5 years alone",
        "ST3 amplification: policy is pouring fuel on the ST1 supply fire",
        "All this pressure lands on corporate balance sheets — Raunak next",
    ]
    add_bullet_box(sl, bullets,
                   Inches(8.8), Inches(2.7), Inches(4.0), Inches(3.5),
                   font_size=15, color=TEXT_DIM)

    speaker_tag(sl, "Tarun → handoff to Raunak")
    return sl


def slide_11_st4_heatmap(prs):
    """Slide 11 — ST4: CapEx Heatmap"""
    sl = blank_slide(prs)
    fill_bg(sl)
    section_bar(sl, "ST4 — Corporate CapEx", ST4_COLOR)

    add_text(sl, "Corporations Are Spending Their Way Out of Risk",
             Inches(0.4), Inches(0.5), Inches(12), Inches(0.7),
             font_size=28, bold=True, color=WHITE)

    add_image(sl, CHARTS / "st4_capex_intensity_heatmap.png",
              Inches(0.4), Inches(1.3), Inches(8.2), Inches(5.5))

    bullets = [
        "10 companies across 4 sectors: semiconductor, hyperscaler, energy, mining",
        "Darkening gradient in semiconductor rows post-2018 = the signal",
        "NVDA, INTC, AMD building redundant capacity — not for efficiency, but resilience",
        "r = 0.56 between annual GPR and semiconductor CapEx intensity",
    ]
    add_bullet_box(sl, bullets,
                   Inches(8.8), Inches(1.5), Inches(4.0), Inches(4.5),
                   font_size=15, color=TEXT_DIM)

    speaker_tag(sl, "Raunak")
    return sl


def slide_12_st4_gpr(prs):
    """Slide 12 — ST4: GPR vs CapEx"""
    sl = blank_slide(prs)
    fill_bg(sl)
    section_bar(sl, "ST4 — Corporate CapEx", ST4_COLOR)

    add_text(sl, "GPR Spikes Predict CapEx Increases 1–2 Years Later",
             Inches(0.4), Inches(0.5), Inches(12), Inches(0.7),
             font_size=28, bold=True, color=WHITE)

    add_image(sl, CHARTS / "st4_gpr_vs_capex.png",
              Inches(0.4), Inches(1.3), Inches(8.0), Inches(5.5))

    add_rect(sl, Inches(8.8), Inches(1.5), Inches(4.0), Inches(1.0), SURFACE)
    add_text(sl, "r = 0.56",
             Inches(8.9), Inches(1.55), Inches(3.8), Inches(0.9),
             font_size=36, bold=True, color=ST4_COLOR)

    bullets = [
        "Intel $20B Ohio fab — announced 2022",
        "TSMC Arizona — announced 2021",
        "Samsung Texas — announced 2021",
        "All right after the GPR spike. Not coincidence — boardrooms pricing risk.",
    ]
    add_bullet_box(sl, bullets,
                   Inches(8.8), Inches(2.7), Inches(4.0), Inches(3.5),
                   font_size=15, color=TEXT_DIM)

    speaker_tag(sl, "Raunak → handoff to Tarun")
    return sl


def slide_13_integration_chain(prs):
    """Slide 13 — Integration: Causal Chain"""
    sl = blank_slide(prs)
    fill_bg(sl)
    section_bar(sl, "Integration", ACCENT_PINK)

    add_text(sl, "The Full Chain — From Shock to Spending",
             Inches(0.4), Inches(0.5), Inches(12), Inches(0.7),
             font_size=28, bold=True, color=WHITE)

    add_image(sl, CHARTS / "integration_causal_chain.png",
              Inches(0.4), Inches(1.3), Inches(8.0), Inches(5.5))

    bullets = [
        "All 4 subtopics on the same time axis",
        "GPR spikes → price vol follows → renewables rise → CapEx climbs",
        "GSCPI → semiconductor CapEx: r = 0.71",
        "Renewables share → semiconductor CapEx: r = 0.82",
        "The chain is real, consistent, and strengthening",
    ]
    add_bullet_box(sl, bullets,
                   Inches(8.8), Inches(1.5), Inches(4.0), Inches(5.0),
                   font_size=15, color=TEXT_DIM)

    speaker_tag(sl, "Tarun")
    return sl


def slide_14_regime(prs):
    """Slide 14 — Integration: Four-Box Regime"""
    sl = blank_slide(prs)
    fill_bg(sl)
    section_bar(sl, "Integration", ACCENT_PINK)

    add_text(sl, "Double Risk Years See Maximum CapEx Deployment",
             Inches(0.4), Inches(0.5), Inches(12), Inches(0.7),
             font_size=28, bold=True, color=WHITE)

    add_image(sl, CHARTS / "integration_four_box_regime.png",
              Inches(0.4), Inches(1.3), Inches(8.0), Inches(5.5))

    bullets = [
        "Each bubble = one year (2010–2022)",
        "Bubble size = semiconductor CapEx intensity",
        "Top-right (Double Risk): 2018–2022 cluster here",
        "Largest bubbles are in the highest-risk quadrant",
        "Firms deploy capital hardest when risks compound",
    ]
    add_bullet_box(sl, bullets,
                   Inches(8.8), Inches(1.5), Inches(4.0), Inches(5.0),
                   font_size=15, color=TEXT_DIM)

    speaker_tag(sl, "Tarun")
    return sl


def slide_15_closing(prs):
    """Slide 15 — Closing"""
    sl = blank_slide(prs)
    fill_bg(sl)

    add_text(sl, "What We Found",
             Inches(0.5), Inches(0.5), Inches(12), Inches(0.8),
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Four-box summary
    findings = [
        (ST1_COLOR, "ST1", "Gallium, Germanium, Rare Earths\nHHI > 0.75 — near-monopolies"),
        (ST2_COLOR, "ST2", "GPR → GSCPI lag: 1–2 months\nTransmission strengthening post-2020"),
        (ST3_COLOR, "ST3", "Solar grew 30× in 12 years\nPolicy is amplifying ST1 demand"),
        (ST4_COLOR, "ST4", "GPR → CapEx: r = 0.56\nFirms price in risk with capital"),
    ]
    for i, (color, tag, text) in enumerate(findings):
        col = i % 2
        row = i // 2
        left = Inches(0.5 + col * 6.2)
        top  = Inches(1.6 + row * 2.3)
        add_rect(sl, left, top, Inches(5.8), Inches(2.0), SURFACE)
        add_rect(sl, left, top, Inches(5.8), Inches(0.07), color)
        add_text(sl, tag, left + Inches(0.15), top + Inches(0.12),
                 Inches(1), Inches(0.4),
                 font_size=14, bold=True, color=color)
        add_text(sl, text, left + Inches(0.15), top + Inches(0.55),
                 Inches(5.5), Inches(1.3),
                 font_size=16, color=TEXT_MAIN)

    add_text(sl,
             "The chain is measurable · consistent across 13 years · getting stronger\n"
             "This is also the analytical foundation for a geopolitical supply chain intelligence SaaS.",
             Inches(0.5), Inches(6.5), Inches(12), Inches(0.8),
             font_size=14, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    speaker_tag(sl, "Tarun")
    return sl


# ══════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════

def build():
    prs = new_prs()

    builders = [
        slide_01_hook,
        slide_02_thesis,
        slide_03_st1_hhi,
        slide_04_st1_sankey,
        slide_05_st1_restrictions,
        slide_06_st2_gpr,
        slide_07_st2_transmission,
        slide_08_st2_ukraine,
        slide_09_st3_mix,
        slide_10_st3_demand,
        slide_11_st4_heatmap,
        slide_12_st4_gpr,
        slide_13_integration_chain,
        slide_14_regime,
        slide_15_closing,
    ]

    for i, builder in enumerate(builders, 1):
        builder(prs)
        print(f"  Built slide {i:02d}: {builder.__name__}")

    out = OUTPUTS / "Vaartha_Presentation.pptx"
    prs.save(str(out))
    print(f"\nSaved: {out}")
    print("Upload to Google Drive → File → Save as Google Slides")


if __name__ == "__main__":
    build()
